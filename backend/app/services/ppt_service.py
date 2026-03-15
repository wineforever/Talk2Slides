import os
import subprocess
import tempfile
import platform
import locale
from pathlib import Path
from typing import List, Dict, Any, Optional
import json
import logging

from pptx import Presentation
from pdf2image import convert_from_path
from pdf2image.exceptions import PDFInfoNotInstalledError
import shutil

from app.core.config import settings

logger = logging.getLogger(__name__)

class PPTService:
    """PPT处理服务"""

    _cached_libreoffice_path: str = ""
    _cached_poppler_path: str = ""
    _poppler_checked: bool = False
    _poppler_available: bool = False

    @staticmethod
    def _safe_text_encoding() -> str:
        enc = locale.getpreferredencoding(False)
        return enc or "utf-8"
    
    def _get_libreoffice_path(self) -> str:
        """获取LibreOffice可执行文件路径
        
        如果配置的路径是相对路径且不可用，尝试在常见位置查找。
        
        Returns:
            LibreOffice可执行文件完整路径
        """
        cls = type(self)
        if cls._cached_libreoffice_path and os.path.exists(cls._cached_libreoffice_path):
            return cls._cached_libreoffice_path

        configured_path = settings.LIBREOFFICE_PATH
        
        # 如果路径已包含目录分隔符，直接返回
        if os.path.sep in configured_path or ('/' in configured_path and os.path.sep == '/'):
            cls._cached_libreoffice_path = configured_path
            return configured_path
        
        # 检查配置的路径是否可直接执行（在PATH中）
        try:
            # 使用shutil.which查找可执行文件
            which_path = shutil.which(configured_path)
            if which_path:
                cls._cached_libreoffice_path = which_path
                return which_path
        except:
            pass
        
        # 如果在Windows上，尝试常见安装位置
        if platform.system() == "Windows":
            common_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                os.path.expandvars(r"%ProgramFiles%\LibreOffice\program\soffice.exe"),
                os.path.expandvars(r"%ProgramW6432%\LibreOffice\program\soffice.exe"),
            ]
            
            for path in common_paths:
                if os.path.exists(path):
                    cls._cached_libreoffice_path = path
                    return path
        
        # 如果都没找到，返回配置的路径（让后续错误处理）
        return configured_path
    
    def _check_poppler_installed(self) -> bool:
        """检查poppler是否安装并可用
        
        Returns:
            True如果poppler已安装并可用，否则False
        """
        cls = type(self)
        if cls._poppler_checked:
            return cls._poppler_available

        poppler_path = self._get_poppler_path()
        executable_name = "pdftoppm.exe" if platform.system() == "Windows" else "pdftoppm"
        command = [str(Path(poppler_path) / executable_name), "-v"] if poppler_path else [executable_name, "-v"]

        try:
            # 尝试运行pdftoppm命令检查poppler是否可用
            result = subprocess.run(
                command,
                capture_output=True,
                text=True,
                encoding=self._safe_text_encoding(),
                errors="replace",
                timeout=5
            )
            cls._poppler_available = (
                result.returncode == 0 or "pdftoppm" in result.stderr or "pdftoppm" in result.stdout
            )
        except (FileNotFoundError, subprocess.TimeoutExpired):
            cls._poppler_available = False

        if cls._poppler_available and poppler_path:
            logger.info("Using Poppler from %s", poppler_path)

        cls._poppler_checked = True
        return cls._poppler_available

    def _get_poppler_path(self) -> Optional[str]:
        cls = type(self)
        if cls._cached_poppler_path and Path(cls._cached_poppler_path).exists():
            return cls._cached_poppler_path

        executable_name = "pdftoppm.exe" if platform.system() == "Windows" else "pdftoppm"
        seen: set[str] = set()

        for candidate in self._iter_poppler_candidates():
            normalized = os.path.normcase(str(candidate))
            if normalized in seen:
                continue
            seen.add(normalized)

            executable_path = candidate / executable_name
            if executable_path.exists():
                cls._cached_poppler_path = str(candidate)
                return cls._cached_poppler_path

        which_path = shutil.which(executable_name)
        if which_path:
            cls._cached_poppler_path = str(Path(which_path).parent)
            return cls._cached_poppler_path

        return None

    def _iter_poppler_candidates(self) -> List[Path]:
        candidates: List[Path] = []

        for configured_path in (
            getattr(settings, "POPPLER_PATH", ""),
            os.getenv("PDF2IMAGE_POPPLER_PATH", ""),
        ):
            candidates.extend(self._expand_poppler_candidate(configured_path))

        if platform.system() == "Windows":
            for root_text in (
                os.environ.get("ProgramFiles", ""),
                os.environ.get("ProgramW6432", ""),
                os.environ.get("ProgramFiles(x86)", ""),
                os.environ.get("LOCALAPPDATA", ""),
            ):
                if not root_text:
                    continue

                root = Path(os.path.expandvars(root_text))
                candidates.extend(
                    [
                        root / "poppler" / "Library" / "bin",
                        root / "poppler" / "bin",
                    ]
                )

                if root.exists():
                    for match in sorted(root.glob("poppler*")):
                        if not match.is_dir():
                            continue
                        candidates.extend(
                            [
                                match / "Library" / "bin",
                                match / "bin",
                            ]
                        )

            candidates.extend(
                [
                    Path(r"C:\poppler\Library\bin"),
                    Path(r"C:\poppler\bin"),
                    Path(r"C:\ProgramData\chocolatey\bin"),
                ]
            )

        return candidates

    def _expand_poppler_candidate(self, configured_path: str) -> List[Path]:
        text = str(configured_path or "").strip().strip('"')
        if not text:
            return []

        candidate = Path(os.path.expandvars(os.path.expanduser(text)))
        if not any(sep in text for sep in ("\\", "/")) and not candidate.drive:
            which_path = shutil.which(text)
            return [Path(which_path).parent] if which_path else []

        if candidate.suffix.lower() == ".exe":
            return [candidate.parent]

        return [
            candidate,
            candidate / "Library" / "bin",
            candidate / "bin",
        ]

    def _build_poppler_error_message(self) -> str:
        configured_path = str(getattr(settings, "POPPLER_PATH", "") or "").strip() or "(not set)"

        if platform.system() == "Windows":
            return (
                "检测到Poppler不可用。pdf2image需要poppler-utils来解析PDF文件。\n\n"
                f"当前POPPLER_PATH配置：{configured_path}\n"
                "如果你已经安装了Poppler，通常是当前服务进程没有继承到正确的PATH。\n\n"
                "请优先在 talk2slides.ini 的 [env] 段显式配置：\n"
                "POPPLER_PATH = C:\\Program Files\\poppler\\Library\\bin\n\n"
                "也可以配置为 pdftoppm.exe 的完整路径，例如：\n"
                "POPPLER_PATH = C:\\Program Files\\poppler\\Library\\bin\\pdftoppm.exe\n\n"
                "如果使用的是带版本号的目录，请改成你的实际目录后重启应用或重启 Windows Service。"
            )

        return (
            "检测到Poppler不可用。pdf2image需要poppler-utils来解析PDF文件。\n\n"
            f"当前POPPLER_PATH配置：{configured_path}\n"
            "请确保 pdftoppm 在 PATH 中可用，或显式配置 POPPLER_PATH 指向 Poppler 的 bin 目录。"
        )
    
    def extract_slides(self, pptx_path: str) -> List[Dict[str, Any]]:
        """从PPTX文件中提取幻灯片文本内容
        
        Args:
            pptx_path: PPTX文件路径
            
        Returns:
            幻灯片列表，每个幻灯片包含标题、正文、备注和索引
        """
        try:
            presentation = Presentation(pptx_path)
            slides = []
            
            for i, slide in enumerate(presentation.slides):
                slide_data = {
                    "index": i,
                    "title": "",
                    "content": "",
                    "notes": "",
                    "full_text": ""
                }
                
                # 提取标题（通常来自标题占位符）
                title_shapes = []
                content_shapes = []
                
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    text = shape.text.strip()
                    if not text:
                        continue
                    
                    # 简单启发式：判断是否为标题
                    # 1. 形状名称包含"title"
                    # 2. 文本较短且字体较大（这里简化处理）
                    # 3. 使用占位符类型
                    if shape.is_placeholder:
                        ph = shape.placeholder
                        if ph.type == 1:  # 标题
                            title_shapes.append(text)
                        else:
                            content_shapes.append(text)
                    else:
                        # 非占位符形状，根据文本长度判断
                        if len(text.split()) <= 10 and len(text) < 100:
                            title_shapes.append(text)
                        else:
                            content_shapes.append(text)
                
                # 提取备注框内容
                notes_text = []
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for shape in notes_slide.shapes:
                        if shape.has_text_frame:
                            text = shape.text.strip()
                            if text:
                                notes_text.append(text)
                
                # 合并标题、内容和备注
                slide_data["title"] = " ".join(title_shapes)
                slide_data["content"] = " ".join(content_shapes)
                slide_data["notes"] = " ".join(notes_text)
                
                # 构建完整文本，备注内容可能很重要，所以包含在内
                full_text_parts = []
                if slide_data["title"]:
                    full_text_parts.append(slide_data["title"])
                if slide_data["content"]:
                    full_text_parts.append(slide_data["content"])
                if slide_data["notes"]:
                    full_text_parts.append(slide_data["notes"])
                
                slide_data["full_text"] = " ".join(full_text_parts).strip()
                
                # 记录提取的备注长度（用于调试）
                if slide_data["notes"]:
                    logger.debug(f"幻灯片 {i}: 提取备注长度 {len(slide_data['notes'])} 字符")
                else:
                    logger.debug(f"幻灯片 {i}: 无备注内容")
                
                slides.append(slide_data)
            
            # 统计备注提取情况
            slides_with_notes = sum(1 for slide in slides if slide["notes"])
            logger.info(f"PPT解析完成: 共 {len(slides)} 张幻灯片，其中 {slides_with_notes} 张包含备注")
            
            return slides
            
        except Exception as e:
            raise Exception(f"PPT解析失败: {str(e)}")
    
    def export_slides_to_images(
        self, 
        pptx_path: str, 
        output_dir: str, 
        resolution: str = "1920x1080"
    ) -> List[str]:
        """将PPTX幻灯片导出为图片
        
        Args:
            pptx_path: PPTX文件路径
            output_dir: 输出目录
            resolution: 输出分辨率，格式为"宽x高"
            
        Returns:
            图片路径列表，按幻灯片顺序排列
        """
        # 在try块之前获取LibreOffice路径，以便在异常处理中使用
        libreoffice_path = self._get_libreoffice_path()
        poppler_path = self._get_poppler_path()
        
        # 检查poppler是否安装
        if not self._check_poppler_installed():
            raise Exception(self._build_poppler_error_message())
        
        try:
            # 确保输出目录存在
            output_path = Path(output_dir)
            output_path.mkdir(parents=True, exist_ok=True)
            
            # 创建临时目录用于PDF转换
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_pdf = Path(temp_dir) / "presentation.pdf"
                
                # 使用LibreOffice将PPTX转换为PDF
                # soffice --headless --convert-to pdf --outdir <dir> <file>
                libreoffice_cmd = [
                    libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", temp_dir,
                    pptx_path
                ]
                
                result = subprocess.run(
                    libreoffice_cmd,
                    capture_output=True,
                    text=True,
                    encoding=self._safe_text_encoding(),
                    errors="replace",
                    timeout=300  # 5分钟超时
                )
                
                if result.returncode != 0:
                    raise Exception(f"PDF转换失败: {result.stderr}")
                
                if not temp_pdf.exists():
                    # 尝试查找生成的PDF文件
                    pdf_files = list(Path(temp_dir).glob("*.pdf"))
                    if not pdf_files:
                        raise Exception("PDF文件未生成")
                    temp_pdf = pdf_files[0]
                
                # 使用pdf2image将PDF转换为图片
                width, height = map(int, resolution.split('x'))
                
                images = convert_from_path(
                    str(temp_pdf),
                    dpi=max(72, int(settings.PPT_EXPORT_DPI)),  # 可通过配置调整导出速度/清晰度
                    size=(width, height),
                    output_folder=str(output_path),
                    fmt="png",
                    paths_only=True,
                    poppler_path=poppler_path,
                )
                
                # 确保图片按顺序命名
                image_paths = []
                for i, image_path in enumerate(sorted(images, key=lambda x: x)):
                    # 重命名为 slide_001.png, slide_002.png 等
                    new_name = output_path / f"slide_{i+1:03d}.png"
                    if image_path != new_name:
                        shutil.move(image_path, new_name)
                    image_paths.append(str(new_name))
                
                return image_paths
                
        except subprocess.TimeoutExpired:
            raise Exception(
                "PDF转换超时。可能原因：\n"
                "1. LibreOffice未正确安装或启动失败\n"
                "2. PPTX文件过大或损坏\n"
                "3. 系统资源不足\n"
                "请检查LibreOffice安装并尝试简化PPTX文件"
            )
        except FileNotFoundError:
            if platform.system() == "Windows":
                raise Exception(
                    f"未找到LibreOffice可执行文件。\n"
                    f"尝试的路径：{libreoffice_path}\n\n"
                    "请执行以下步骤：\n"
                    "1. 下载安装LibreOffice：https://www.libreoffice.org/download/download/\n"
                    "2. 安装时选择'添加到PATH'选项\n"
                    "3. 或者设置LIBREOFFICE_PATH环境变量指向soffice.exe完整路径\n"
                    f"当前配置路径：{settings.LIBREOFFICE_PATH}\n"
                    f"常见安装路径：C:\\Program Files\\LibreOffice\\program\\soffice.exe"
                )
            else:
                raise Exception(
                    f"未找到LibreOffice可执行文件。\n"
                    f"尝试的路径：{libreoffice_path}\n\n"
                    "请执行以下步骤：\n"
                    "1. Linux安装：sudo apt-get install libreoffice 或 sudo yum install libreoffice\n"
                    "2. macOS安装：brew install libreoffice\n"
                    "3. 或者设置LIBREOFFICE_PATH环境变量指向soffice完整路径\n"
                    f"当前配置路径：{settings.LIBREOFFICE_PATH}"
                )
        except PDFInfoNotInstalledError:
            raise Exception(self._build_poppler_error_message())
        except Exception as e:
            raise Exception(f"PPT导出失败: {str(e)}")
    
    def get_slide_count(self, pptx_path: str) -> int:
        """获取幻灯片数量"""
        try:
            presentation = Presentation(pptx_path)
            return len(presentation.slides)
        except Exception as e:
            raise Exception(f"获取幻灯片数量失败: {str(e)}")
    
    def validate_pptx(self, pptx_path: str) -> bool:
        """验证PPTX文件是否有效"""
        try:
            presentation = Presentation(pptx_path)
            return len(presentation.slides) > 0
        except:
            return False
